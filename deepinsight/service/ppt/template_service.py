from __future__ import annotations

import io
import json
import os
from collections import deque
from typing import Dict, Any, Tuple, List, Literal
from pathlib import Path
import json
import logging
import re
from copy import deepcopy
import csv

from bs4 import BeautifulSoup
from bs4.element import NavigableString, Tag
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt, Cm

from deepinsight.utils.md_render import to_html

SLIDE_INDEX_MAP = {
    "cover_page": 0,
    "content_page": 1,
    "conf_overview_page": 2,
    "tech_theme_page": 3,
    "research_hotspot_collab_01_page": 4,
    "research_hotspot_collab_02_page": 5,
    "country_tech_feature_page": 6,
    "institution_tech_feature_page": 7,
    "institution_tech_strength_page": 8,
    "institution_cooperation_page": 9,
    "high_potential_tech_transfer_page": 10,
    "keynote_page": 11,
    "topic_content_page": 12,
    "topic_detail_page": 13,
    "valuable_paper_page": 14,
    "conf_summary_page": 15,
}

TABLE_DEFAULT_STYLE={
    "font_name": "微软雅黑",
    "font_size": 8,
    "bold": False,
    "color": [0, 0, 0],
    "align": PP_ALIGN.LEFT,
    "bg_color": [255, 255, 255],

    # 表头样式
    "header_font_name": "微软雅黑",
    "header_font_size": 9,
    "header_bold": True,
    "header_color": [0, 0, 0],
    "header_bg_color": [233, 233, 233],  # 蓝色背景
    "header_align": PP_ALIGN.LEFT,

    # 尺寸控制
    "row_height": None,
    "first_col_width": 1.2
}

topic_table_style = deepcopy(TABLE_DEFAULT_STYLE)
topic_table_style["first_col_width"] = 2
topic_table_style["font_size"] = 11
topic_table_style["header_font_size"] = 13
high_potential_table_stype = deepcopy(TABLE_DEFAULT_STYLE)
high_potential_table_stype["first_col_width"] = 3.5
institution_tech_table_stype = deepcopy(TABLE_DEFAULT_STYLE)
institution_tech_table_stype["first_col_width"] = 2.5
hotpots_table_stype = deepcopy(TABLE_DEFAULT_STYLE)
hotpots_table_stype["font_size"] = 10

STYLE_CONFIG = {
    "cover_page": {
        "CONFERENCE_NAME": {"font_size": 44, "bold": True, "color": [192, 0, 0], "align": PP_ALIGN.LEFT},
        "DATE": {"font_size": 24, "bold": False, "color": [0, 0, 0], "align": PP_ALIGN.LEFT},
    },
    "content_page": {},
    "conf_overview_page": {
        "conf_info": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "organizer_level": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "conf_topics": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "conf_loc": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "conf_date": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "conf_sponsor": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "conf_chair": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "conf_committee": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "conf_institution": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "submit_papers": {"font_size": 12, "bold": False, "color": [50, 50, 50]},
        "total_trend": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
    },

    # === 新增：技术主题分析 ===
    "tech_theme_page": {
        "tech_field_png": {},
        "key_tech_intro": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "key_tech_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 新增：研究热点与跨区域技术合作 01 ===
    "research_hotspot_collab_01_page": {
        "keyword_cloud_png": {},
        "keyword_intro": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keyword_couple_analysis_png": {},
        "keyword_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 新增：研究热点与跨区域技术合作 02 ===
    "research_hotspot_collab_02_page": {
        "keyword_topic_csv": hotpots_table_stype,  # 重用已有的表格样式或按需替换
        "keyword_topic_intro": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keyword_topic_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 新增：国家/地区技术特征分析 ===
    "country_tech_feature_page": {
        "country_tech_top_png": {},
        "country_tech_strength_csv": TABLE_DEFAULT_STYLE,
        "country_tech_intro": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "country_tech_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 新增：机构技术特征分析 ===
    "institution_tech_feature_page": {
        "top_institution_png": {},
        "institution_tech_feat_intro": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "company_school_analysis_png": {},
        "institution_tech_feat_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 新增：机构技术优势分析 ===
    "institution_tech_strength_page": {
        "university_tech_strength_csv": institution_tech_table_stype,
        "company_tech_strength_csv": institution_tech_table_stype,
        "institution_tech_strength_intro": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "institution_tech_strength_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 新增：跨机构合作网络分析 ===
    "institution_cooperation_page": {
        "institution_cooperation_png": {},
        "institution_cooperation_intro": {"font_size": 11, "bold": False, "color": [0, 0, 0]},
        "institution_cooperation_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 新增：高潜技术转化分析 ===
    "high_potential_tech_transfer_page": {
        "high_potential_csv": high_potential_table_stype,
        "high_potential_intro": {"font_size": 11, "bold": False, "color": [0, 0, 0]},
        "high_potential_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
    },

    # === 后续原有页面（序号顺延） ===
    "keynote_page": {
        "keynote_title": {"font_size": 28, "bold": True, "color": [153, 0, 0]},
        "speaker": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keynote_abstract": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keynote_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
        "keynote_background": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keynote_objective": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keynote_method": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keynote_inspiration": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "keynote_pic1_png": {},
        "keynote_pic2_png": {},
    },
    "topic_content_page": {
        "topic_content_csv": topic_table_style
    },
    "topic_detail_page": {
        "topic_title": {"font_size": 24, "bold": True, "color": [153, 0, 0]},
        "topic_overview": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "topic_reason": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "topic_summary":  {"font_size": 12, "bold": False, "color": [255, 255, 255]},
        "topic_method_innovation": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "topic_inspiration": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
    },
    "valuable_paper_page": {
        "tech_topic": {"font_size": 11, "bold": True, "color": [255, 255, 255]},
        "paper_title": {"font_size": 9, "bold": False, "color": [0, 0, 0]},
        "paper_headline": {"font_size": 24, "bold": True, "color": [153, 0, 0]},
        "paper_background": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "paper_challenges": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "paper_tech_resource": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "paper_key_tech": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "paper_result": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "paper_summary": {"font_size": 12, "bold": False, "color": [255, 255, 255]},
        "key_tech_png": {},
        "exp_result_png": {},
    },
    "conf_summary_page": {
        "key_trends": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
        "suggestions": {"font_size": 12, "bold": False, "color": [0, 0, 0]},
    }
}

PRE_STR_MAP = {
    "keynote_page": "",
    "topic_detail_page": "",
    "valuable_paper_page": ""
}

class PPTTemplateService:
    """
    根据 JSON 数据填充 PPT 模板的服务。

    使用规则：
    - 文本占位符：形状或表格单元格中的文本包含 `{{key}}` 时，替换为 JSON 中的对应值。
    - 暂不处理图片与图表的动态替换，后续可扩展。
    """
    @staticmethod
    def _cleanup_html_breaks(cell: Tag, keep_bold_mark: bool, bold_mark: Literal["**", "__"]) -> str:
        skip = {"code", "pre", "span", "br"}
        bold = {"b", "strong"}
        queue = deque(cell.contents)
        while queue:
            node = queue.popleft()
            if node.name in bold and keep_bold_mark:  # noqa: bs do not make a typehint on .name
                node.insert_before(NavigableString(bold_mark))
                node.insert_after(NavigableString(bold_mark))
            if isinstance(node, NavigableString):
                node.replace_with(re.sub(r"\s+", " ", node))
            elif node.name in skip:  # noqa: bs do not make a typehint on .name
                continue
            elif hasattr(node, "contents"):
                queue.extend(node.contents)

        for br in cell.find_all("br"):
            br.replace_with("\n")

        return cell.get_text(strip=False)

    @classmethod
    def load_md_table(cls, md: str, keep_bold_mark: bool = True, bold_mark: Literal["**", "__"] = "**") -> list[list[str]]:
        html = to_html(md)[0]
        soup = BeautifulSoup(html, "lxml")
        html_tables = soup.find_all("table")
        if len(html_tables) > 1:
            logging.warning("More than one table found. Only the first table is reserved.")
        elif not html_tables:
            logging.warning("No table found.")
            return []
        html_table = html_tables[0]
        out_table = []
        for row in html_table.find_all("tr"):
            cells = row.find_all(['th', 'td'])
            out_row = []
            for cell in cells:
                out_row.append(cls._cleanup_html_breaks(cell, keep_bold_mark, bold_mark))
            out_table.append(out_row)
        return out_table

    def fill_from_json_file(self, template_path: str, json_file_path: str, output_name: str | None = None) -> Presentation:
        """
        从 JSON 文件读取数据并填充 PPT 模板。

        Args:
            template_path: PPT 模板文件路径 (.pptx)
            json_file_path: 内容 JSON 文件路径
            output_name: 输出文件名（可选），默认 `result.pptx`

        Returns:
            Presentation: 填充后的 Presentation 对象
        """
        with open(json_file_path, "r", encoding="utf-8") as f:
            list_json = json.load(f)
        return self.fill_from_json(template_path, list_json, output_name=output_name)

    def fill_from_json(self, template_path: str, list_json: Dict[str, Any], output_name: str | None = None) -> Presentation:
        """
        使用内存中的 JSON 数据填充 PPT 模板。

        Args:
            template_path: PPT 模板文件路径 (.pptx)
            data: 用于填充的字典数据
            output_name: 输出文件名（可选），默认 `result.pptx`

        Returns:
            Presentation: 填充后的 Presentation 对象
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"PPT 模板不存在: {template_path}")

        pres = Presentation(template_path)
        original_count = len(pres.slides)

        for item in list_json:
            t = item.get("type")
            content = item.get("content")
            skip_fill = item.get("skip_fill", False)
            if t not in SLIDE_INDEX_MAP:
                continue
            template_slide = pres.slides[SLIDE_INDEX_MAP[t]]
            new_slide = self.duplicate_slide_to_end(template_slide, pres)
            if content and not skip_fill:
                self.replace_content(new_slide, t, content)

        self.delete_slides(pres, list(range(original_count)))
        return pres
    
    def duplicate_slide_to_end(self, template_slide, pres):
        new_slide = pres.slides.add_slide(template_slide.slide_layout)
        for shape in list(new_slide.shapes):
            shape.element.getparent().remove(shape.element)
        for shape in template_slide.shapes:
            new_el = deepcopy(shape.element)
            new_slide.shapes._spTree.append(new_el)
        if template_slide.background:
            try:
                new_slide.background._element = deepcopy(template_slide.background._element)
            except:
                pass
        return new_slide

    def replace_content(self, slide, type_name: str, data: Any):
        if not isinstance(data, dict):
            return
        placeholders = {self._shape_text(s): s for s in slide.shapes if self._shape_text(s)}
        for key, val in data.items():
            if key not in placeholders:
                continue
            shape = placeholders[key]
            conf = STYLE_CONFIG.get(type_name, {}).get(key, {})
            if isinstance(val, dict) and val.get("type") == "image":
                path = val["path"]
                if path and Path(path).exists():
                    self._insert_image(slide, shape, path, conf)
            elif (isinstance(val, dict) and val.get("type") == "table") or (
                    isinstance(val, str) and val.endswith(".csv")):
                path = val.get("path") if isinstance(val, dict) else val
                conf_table = STYLE_CONFIG.get(type_name, {}).get(key, {})
                md_str = val.get("content") if isinstance(val, dict) else val
                if path and Path(path).exists():
                    self._insert_table_from_csv(slide, shape, csv_path=path, conf=conf_table)
                elif md_str:
                    self._insert_table_from_md(slide, shape, md=md_str, conf=conf_table)

            else:
                text = val["text"] if isinstance(val, dict) and val.get("text") else val
                pre_str = PRE_STR_MAP.get(type_name, "")
                if not text:
                    text = ""
                if "TITLE" in key:
                    text = pre_str + text
                self._insert_text_from_markdown(shape, text, conf)

    def _shape_text(self, shape) -> str:
        return getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""


    def _replace_text(self, shape, text: str, conf: dict):
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = str(text)
            font_size = Pt(conf.get("font_size")) if conf.get("font_size") else None
            bold = conf.get("bold")
            font_color = RGBColor(*conf["color"]) if conf.get("color") else None
            align = conf.get("align")
            font_name = conf.get("font_name", "Microsoft YaHei")  # 默认微软雅黑
            if font_size or bold is not None or font_color:
                for run in p.runs:
                    if font_name:
                        run.font.name = font_name
                    if font_size:
                        run.font.size = font_size
                    if bold is not None:
                        run.font.bold = bold
                    if font_color:
                        run.font.color.rgb = font_color
            if align:
                p.alignment = align


    def _insert_image(self, slide, shape, image_path: str, conf: dict):
        left, top = shape.left, shape.top
        width = Inches(conf.get("width")) if conf.get("width") else shape.width
        height = Inches(conf.get("height")) if conf.get("height") else shape.height
        width = Cm(conf.get("width_cm")) if conf.get("width_cm") else width
        height = Cm(conf.get("height_cm")) if conf.get("height_cm") else height
        try:
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        except Exception as e:
            self._replace_text(shape, "", conf)
            return
        slide.shapes._spTree.remove(shape.element)  # 删除模板图占位


    def _add_solid_fill_to_tcPr(self, tcPr, rgb_tuple):
        """
        在单元格 tcPr 内附加一个 <a:solidFill><a:srgbClr val="RRGGBB"/></a:solidFill>
        使用 OxmlElement，兼容不同 pptx 版本。
        """
        # create a:solidFill
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        # set hex color value
        hexval = "%02X%02X%02X" % tuple(rgb_tuple)
        srgbClr.set('val', hexval)
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    def _insert_table_from_md(self, slide, template_shape, md: str, conf: dict) -> None:
        table = self.load_md_table(md, keep_bold_mark=True)
        self._insert_table_from_array(slide, template_shape, table, conf)

    def _insert_table_from_csv(self, slide, template_shape, csv_path: str = None, csv_str: str = None, conf: dict = None):
        conf = conf or {}
        rows = []
        if csv_path:
            with open(csv_path, newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(row)

        elif csv_str:
            csv_content = io.StringIO(csv_str)
            reader = csv.reader(csv_content)
            for row in reader:
                rows.append(row)
        self._insert_table_from_array(slide, template_shape, rows, conf)

    @staticmethod
    def _insert_table_from_array(slide, template_shape, rows: list[list[str]], conf: dict):
        if not rows:
            return

        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)

        left, top = template_shape.left, template_shape.top

        # sizes
        width_in = conf.get("width", None)
        height_in = conf.get("height", None)
        total_width = Inches(width_in) if width_in else template_shape.width
        total_height = Inches(height_in) if height_in else template_shape.height

        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_width, total_height)
        table = table_shape.table

        # styles
        font_name = conf.get("font_name", "Microsoft YaHei")
        font_size = Pt(conf.get("font_size", 11)) if conf.get("font_size") else Pt(11)
        bold = conf.get("bold", False)
        color = conf.get("color", [0,0,0])
        align = conf.get("align", PP_ALIGN.CENTER)

        header_font_name = conf.get("header_font_name", font_name)
        header_font_size = Pt(conf.get("header_font_size", int(font_size.pt)+1)) if conf.get("header_font_size") else Pt(int(font_size.pt)+1)
        header_bold = conf.get("header_bold", True)
        header_color = conf.get("header_color", [255,255,255])
        header_align = conf.get("header_align", PP_ALIGN.CENTER)
        header_bg_color = conf.get("header_bg_color", None)

        row_height_in = conf.get("row_height", None)  # inches
        first_col_width_in = conf.get("first_col_width", None)  # inches

        # distribute column widths (same logic as your last version)
        if n_cols == 1:
            try:
                table.columns[0].width = total_width
            except Exception:
                pass
        else:
            if first_col_width_in:
                first_col_emu = Inches(first_col_width_in)
                remaining = total_width - first_col_emu
                if remaining <= 0:
                    per_col = total_width // n_cols
                    for j in range(n_cols):
                        try:
                            table.columns[j].width = per_col
                        except Exception:
                            pass
                else:
                    try:
                        table.columns[0].width = first_col_emu
                    except Exception:
                        pass
                    per_col = remaining // (n_cols - 1)
                    for j in range(1, n_cols):
                        try:
                            table.columns[j].width = per_col
                        except Exception:
                            pass
            else:
                per_col = total_width // n_cols
                for j in range(n_cols):
                    try:
                        table.columns[j].width = per_col
                    except Exception:
                        pass

        # set row heights if provided
        if row_height_in:
            try:
                for r in table.rows:
                    r.height = Inches(row_height_in)
            except Exception:
                pass

        # helper: create paragraphs + runs from markdown-like text
        bold_pattern = re.compile(r'\*\*(.+?)\*\*')

        def _fill_cell_with_markdown(cell, text, is_header=False):
            """
            text: may contain '\n' for multiple paragraphs, and **bold** markers.
            Apply header or normal styles accordingly.
            """
            # clear existing text_frame
            tf = cell.text_frame
            tf.clear()

            lines = text.splitlines()

            for idx, line in enumerate(lines):
                if "\\n" in line:
                    line=line.replace("\\n", "\n")
                if idx == 0:
                    para = tf.paragraphs[0]
                    para.clear()
                else:
                    para = tf.add_paragraph()
                    para.clear()

                # alignment
                try:
                    para.alignment = header_align if is_header else align
                except Exception:
                    pass

                last = 0
                any_bold = False
                for m in bold_pattern.finditer(line):
                    any_bold = True
                    s, e = m.span()
                    # preceding
                    if s > last:
                        seg = line[last:s]
                        if seg:
                            run = para.add_run()
                            run.text = seg
                            # normal run style
                            if not is_header:
                                if font_name: run.font.name = font_name
                                run.font.size = font_size
                                run.font.bold = bold
                                run.font.color.rgb = RGBColor(*color)
                            else:
                                if header_font_name: run.font.name = header_font_name
                                run.font.size = header_font_size
                                run.font.bold = header_bold
                                run.font.color.rgb = RGBColor(*header_color)
                    # bold segment
                    inner = m.group(1)
                    run = para.add_run()
                    run.text = inner
                    if not is_header:
                        if font_name: run.font.name = font_name
                        run.font.size = font_size
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(*color)
                    else:
                        if header_font_name: run.font.name = header_font_name
                        run.font.size = header_font_size
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(*header_color)
                    last = e
                # trailing
                if last < len(line):
                    trailing = line[last:]
                    if trailing:
                        run = para.add_run()
                        run.text = trailing
                        if not is_header:
                            if font_name: run.font.name = font_name
                            run.font.size = font_size
                            run.font.bold = bold
                            run.font.color.rgb = RGBColor(*color)
                        else:
                            if header_font_name: run.font.name = header_font_name
                            run.font.size = header_font_size
                            run.font.bold = header_bold
                            run.font.color.rgb = RGBColor(*header_color)

                # if no bold at all and no runs created (unlikely), write whole line
                if not any_bold and not para.runs:
                    run = para.add_run()
                    run.text = line
                    if not is_header:
                        if font_name: run.font.name = font_name
                        run.font.size = font_size
                        run.font.bold = bold
                        run.font.color.rgb = RGBColor(*color)
                    else:
                        if header_font_name: run.font.name = header_font_name
                        run.font.size = header_font_size
                        run.font.bold = header_bold
                        run.font.color.rgb = RGBColor(*header_color)

        # populate table
        for i, row_vals in enumerate(rows):
            is_header = (i == 0)
            for j in range(n_cols):
                val = row_vals[j] if j < len(row_vals) else ""
                cell = table.cell(i, j)
                # if val is not string, coerce
                cell_text = "" if val is None else str(val)
                _fill_cell_with_markdown(cell, cell_text, is_header=is_header)

                # header background if requested
                if is_header and header_bg_color:
                    try:
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        # add solidFill
                        solidFill = OxmlElement('a:solidFill')
                        srgbClr = OxmlElement('a:srgbClr')
                        hexval = "%02X%02X%02X" % tuple(header_bg_color)
                        srgbClr.set('val', hexval)
                        solidFill.append(srgbClr)
                        tcPr.append(solidFill)
                    except Exception:
                        pass

        # remove template placeholder shape (best-effort)
        try:
            slide.shapes._spTree.remove(template_shape.element)
        except Exception:
            pass

    def _insert_text_from_markdown(self, shape, value, conf: dict):
        """
        Robust Markdown-lite -> PPT paragraphs renderer.
        Fixes invisible-char-before-* issue by normalizing lines before detecting bullets.
        Supports:
        ◦ bullet lines starting with '* ', '- ', '+ ' (real bullet via para.level=0)
        ◦ inline bold anywhere (including inside bullets) using **bold**
        ◦ newline -> new paragraph
        ◦ color segments using <color (R,G,B)>...</color> or <color (R,G,B)>...<color/> (both accepted)

        conf supports: font_name, font_size (pt), bold (default), color [R,G,B], align (PP_ALIGN)
        """
        if not hasattr(shape, "text_frame"):
            return

        raw_text = "" if value is None else str(value)
        conf = conf or {}

        font_name = conf.get("font_name", "Microsoft YaHei")
        font_size = Pt(conf.get("font_size")) if conf.get("font_size") else None
        default_bold = conf.get("bold", None)
        conf_color = conf.get("color", None)
        align = conf.get("align", None)

        tf = shape.text_frame
        tf.clear()

        lines = raw_text.splitlines()

        def _create_run(paragraph, txt, make_bold=None, color_rgb=None):
            if not txt:
                return None
            run = paragraph.add_run()
            run.text = txt
            if font_name:
                try:
                    run.font.name = font_name
                except Exception:
                    pass
            if font_size:
                run.font.size = font_size
            # bold precedence: explicit make_bold (True/False) > default_bold if provided
            if make_bold is not None:
                run.font.bold = make_bold
            elif default_bold is not None:
                run.font.bold = default_bold
            # color: color_rgb (tuple) overrides conf_color
            if color_rgb:
                try:
                    run.font.color.rgb = RGBColor(*color_rgb)
                except Exception:
                    pass
            elif conf_color:
                try:
                    run.font.color.rgb = RGBColor(*conf_color)
                except Exception:
                    pass
            return run

        # helper: split a line into segments [(color_rgb_or_None, text), ...]
        def _split_color_segments(line):
            segments = []
            # opening tag: <color (R,G,B)>
            open_pat = re.compile(r'<color\s*\(\s*(\d{1,3})\s*,\s*(\d{1,3})\s*,\s*(\d{1,3})\s*\)\s*>',
                                re.IGNORECASE)
            # closing tag can be </color> or <color/> (allow spaces and case-insensitive)
            close_pat = re.compile(r'(?:</\s*color\s*>|<\s*color\s*/\s*>)', re.IGNORECASE)
            pos = 0
            while pos < len(line):
                m = open_pat.search(line, pos)
                if not m:
                    remainder = line[pos:]
                    if remainder:
                        segments.append((None, remainder))
                    break
                start, end = m.span()
                # text before opening tag
                if start > pos:
                    before = line[pos:start]
                    if before:
                        segments.append((None, before))
                # parse rgb
                try:
                    r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
                    # clamp 0..255
                    r = max(0, min(255, r)); g = max(0, min(255, g)); b = max(0, min(255, b))
                    rgb = (r, g, b)
                except Exception:
                    rgb = None
                # find closing tag after the opening tag
                close_m = close_pat.search(line, end)
                if not close_m:
                    # no close tag: take rest of line
                    content = line[end:]
                    segments.append((rgb, content))
                    break
                else:
                    cstart, cend = close_m.span()
                    content = line[end:cstart]
                    segments.append((rgb, content))
                    pos = cend
                    # continue scanning after close
            return segments

        first_para = True
        bold_pat = re.compile(r'\*\*(.+?)\*\*')

        for orig_line in lines:
            # Normalize leading characters that commonly break ^\s*\* detection:
            line = orig_line.lstrip('\ufeff\u200b\u00A0 \t"\'\u201c\u201d\u2018\u2019')
            if not line.strip():
                continue
            # After stripping control/invisible chars, detect bullet at line start
            is_bullet = bool(re.match(r'^[\*\-\+]\s+', line))

            # If bullet detected, remove first marker and following spaces
            if is_bullet:
                line = re.sub(r'^[\*\-\+]\s+', '', line, count=1)

            # Get paragraph: first exists by default
            if first_para:
                para = tf.paragraphs[0]
                para.clear()
                first_para = False
            else:
                para = tf.add_paragraph()
                para.clear()

            # set alignment if provided
            if align:
                try:
                    para.alignment = align
                except Exception:
                    pass

            # enable real bullet (relies on template placeholder having bullet style)
            if is_bullet:
                try:
                    para.level = 0
                except Exception:
                    pass

            # split into color segments
            segments = _split_color_segments(line)

            # For each color segment, further split into bold/non-bold parts and create runs
            for seg_color, seg_text in segments:
                if not seg_text:
                    continue
                last_idx = 0
                for m in bold_pat.finditer(seg_text):
                    s, e = m.span()
                    # text before bold
                    if s > last_idx:
                        pre = seg_text[last_idx:s]
                        if pre:
                            _create_run(para, pre, make_bold=None, color_rgb=seg_color)
                    # bold text (inner)
                    inner = m.group(1)
                    _create_run(para, inner, make_bold=True, color_rgb=seg_color)
                    last_idx = e
                if last_idx < len(seg_text):
                    trailing = seg_text[last_idx:]
                    if trailing:
                        _create_run(para, trailing, make_bold=None, color_rgb=seg_color)

            # Ensure at least one run exists
            if not para.runs:
                _create_run(para, line, make_bold=None, color_rgb=None)

    def delete_slides(self, pres, indices: List[int]):
        sldIdLst = pres.slides._sldIdLst
        for offset, idx in enumerate(indices):
            rId = sldIdLst[idx - offset]
            sldIdLst.remove(rId)


